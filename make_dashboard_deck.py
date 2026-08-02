# -*- coding: utf-8 -*-
"""Rebuild ActivityDashboard.pptx and the ActivityDashboard-*.png dashboards behind it.

    pip install pandas matplotlib python-pptx
    python make_dashboard_deck.py

Slides: an all-PC overview, week by week, one hour-by-hour slide per week in WEEKS,
then one slide per machine.

Sources, in this order:
  1. activity_slim.csv.gz from the private ActivityLogger-data repo - the work PCs,
     already deduped, window titles stripped. Cloned next to this repo, or point
     ACTIVITYLOG_DATA at the .csv.gz.
  2. This PC's own *_ActivityLog*.csv in the OneDrive Documents\\ActivityLogger folder,
     for any machine the export does not already cover. The folder can hold an
     overlapping OneDrive snapshot copy, so those rows are deduped on
     PC + start + end + window + process; the slim export must NOT be deduped again.

Counting. One person uses all four PCs, so the all-PC slides count each second once:
a vncviewer.exe window is always a view of one of these same machines and earns no time
at all - the work counts on the PC being driven - and any remaining second two PCs both
logged goes to the row that started first. Rows over 2 h are dropped, not clamped: they
are idle-detection failures where a scheduled task woke the PC and the logger kept
crediting the last foreground window. Per-PC slides show each machine's own log in full.
"""
import copy
import glob
import os

import numpy as np
import pandas as pd
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.colors import LinearSegmentedColormap, to_rgb
from matplotlib.patches import FancyBboxPatch
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

REPO = os.path.dirname(os.path.abspath(__file__))
DECK = os.path.join(REPO, "ActivityDashboard.pptx")
SLIM = os.environ.get("ACTIVITYLOG_DATA") or os.path.join(
    os.path.dirname(REPO), "ActivityLogger-data", "data", "activity_slim.csv.gz")

def local_log_dir():
    """The ActivityLogger output folder, found the same way the app picks it."""
    home = os.path.expanduser("~")
    for base in ("OneDrive", "OneDrive - Personal", "OneDrive - GE HealthCare"):
        d = os.path.join(home, base, "Documents", "ActivityLogger")
        if glob.glob(os.path.join(d, "*_ActivityLog*.csv")):
            return d
    d = os.path.join(os.environ.get("LOCALAPPDATA", home), "ActivityLogger")
    return d if glob.glob(os.path.join(d, "*_ActivityLog*.csv")) else None

BG, CARD, LINE = "#F9F9F7", "#FCFCFB", "#E6E5E1"
INK, INK2, MUTED = "#1B1A18", "#55534E", "#898781"
# categorical hues, fixed per PC (never by rank). All pairs clear the CVD check;
# #C7359B was picked over violet/teal/ochre, which collide with blue or orange.
PCS = [("HC-D2BYW34", "Work Dell Desktop", "#2A78D6"),
       ("HC-20NXD54", "W11 GEHC Work Laptop", "#EB6834"),
       ("HC-41BYW34", "MatLab desktop", "#1BAF7A"),
       ("DESKTOP-GHT99IR", "home PC", "#C7359B")]
COLOR = {p: c for p, _, c in PCS}
NAME = {p: n for p, n, _ in PCS}
BLUE_RAMP = LinearSegmentedColormap.from_list(
    "b", ["#EAF1FB", "#96C0F2", "#2A78D6", "#104180", "#0D376D"])
plt.rcParams.update({"font.family": "Segoe UI", "axes.facecolor": CARD, "figure.facecolor": BG})

def mix(c, other, t):
    a, b = np.array(to_rgb(c)), np.array(to_rgb(other))
    return tuple(a * (1 - t) + b * t)

def ramp_for(hue):
    return LinearSegmentedColormap.from_list(
        "r", [mix(hue, "white", .93), mix(hue, "white", .55), hue,
              mix(hue, "black", .35), mix(hue, "black", .5)])

VIEWER = "vncviewer.exe"

# Category hues for the weekly stack. Fixed per category name, not per rank, and capped
# at six: every pair clears the CVD check (worst normal 18.8, worst CVD 9.9), which an
# eight-hue stack could not - blue vs violet collapses to 1.7 under deuteranopia.
CAT_COLORS = [("Web Browsing", "#2A78D6"), ("Development", "#EB6834"),
              ("Email", "#1BAF7A"), ("Uncategorized", "#C7359B"),
              ("Meetings", "#8A5A2B"), ("Other categories", "#C9C7C1")]
CAT_ORDER = [c for c, _ in CAT_COLORS]

# ------------------------------------------------------------------ load
COLS = ["StartTime", "EndTime", "DurationSeconds", "ProcessName",
        "Category", "ComputerName", "WindowTitle"]

def load():
    if not os.path.exists(SLIM):
        raise SystemExit("no activity_slim.csv.gz at %s - clone coffeng/ActivityLogger-data "
                         "next to this repo, or set ACTIVITYLOG_DATA" % SLIM)
    work = pd.read_csv(SLIM, parse_dates=["StartTime", "EndTime"])
    work["WindowTitle"] = ""
    n_local_raw = n_local = 0
    frames = [work[COLS]]

    d = local_log_dir()
    if d:
        files = sorted(glob.glob(os.path.join(d, "*_ActivityLog*.csv")))
        loc = pd.concat([pd.read_csv(f, encoding="utf-8-sig") for f in files],
                        ignore_index=True)
        # run this on a PC the export already covers and its rows would be counted twice
        loc = loc[~loc.ComputerName.isin(work.ComputerName.unique())]
        if len(loc):
            n_local_raw = len(loc)
            loc["StartTime"] = pd.to_datetime(loc.StartTime)
            loc["EndTime"] = pd.to_datetime(loc.EndTime)
            # this folder can hold an overlapping OneDrive snapshot copy; the slim export
            # is already deduped and must NOT be deduped again (window titles were dropped)
            loc = loc.drop_duplicates(subset=["ComputerName", "StartTime", "EndTime",
                                              "WindowTitle", "ProcessName"])
            n_local = len(loc)
            frames.append(loc[COLS])
            print("local logs: %s (%s raw rows -> %s)"
                  % (d, format(n_local_raw, ","), format(n_local, ",")))
    if n_local == 0:
        print("no local logs to add - the export already covers every PC found")
    return pd.concat(frames, ignore_index=True), n_local_raw, n_local, len(work)

FAILED = 7200          # a foreground window held for 2 h+ means idle detection failed

def active(df):
    """Active rows only: no idle, no lock screen, and no idle-detection failures.

    A row longer than 2 h is not someone working - a scheduled task wakes the PC around
    01:38 and the logger keeps crediting whatever window was last in front. Those rows are
    dropped outright rather than clamped, so they cannot pad totals or invent a night shift.
    """
    idle = (df.Category == "Inactive") | (df.ProcessName.str.lower() == "lockapp.exe")
    idle_h = df[idle].DurationSeconds.clip(upper=FAILED).sum() / 3600
    a = df[~idle].copy()
    failed = a.DurationSeconds > FAILED
    n_failed, failed_h = int(failed.sum()), a.loc[failed, "DurationSeconds"].clip(
        upper=FAILED).sum() / 3600
    a = a[~failed].copy()
    a["sec"] = a.DurationSeconds
    a["hours"] = a.sec / 3600
    a["end_c"] = a.StartTime + pd.to_timedelta(a.sec, unit="s")
    a["viewer"] = a.ProcessName.str.lower() == VIEWER
    return a, idle_h, n_failed, failed_h

def to_secs(df):
    return np.stack([df[c].values.astype("datetime64[s]").astype(np.int64)
                     for c in ("StartTime", "end_c")], axis=1)

def merge_intervals(df):
    out = []
    for s, e in to_secs(df.sort_values("StartTime")):
        if out and s <= out[-1][1]:
            out[-1][1] = max(out[-1][1], e)
        else:
            out.append([s, e])
    return np.array(out) if out else np.zeros((0, 2), dtype=np.int64)

def overlap_with(M, s, e):
    """Seconds of [s, e) covered by the merged interval array M."""
    if not len(M):
        return 0
    starts, ends = M[:, 0], M[:, 1]
    i = np.searchsorted(ends, s, "right")
    tot = 0
    while i < len(starts) and starts[i] < e:
        tot += min(e, ends[i]) - max(s, starts[i])
        i += 1
    return tot

def credited_seconds(a):
    """Give every second of wall clock to exactly one row.

    A vncviewer window was always a view of one of the other three PCs, so it never earns
    time of its own: the work it shows belongs to the machine being driven, whether or not
    that machine's logger happened to be running. Everything else is credited to the row
    that started first - one person drives all four PCs, so two of them active in the same
    second is double counting. Sorting by start makes the already-claimed part of any row
    contiguous, so a running cursor is exact.
    """
    credit = pd.Series(0.0, index=a.index)           # viewer rows keep 0
    local = a[~a.viewer].sort_values("StartTime")
    cursor = -np.inf
    for ix, (s, e) in zip(local.index, to_secs(local)):
        credit[ix] = max(0, e - max(s, cursor))
        cursor = max(cursor, e)
    return credit.values

def hour_split(df, col="net"):
    """Weekday x hour matrix, splitting each row across the clock hours it really covers.

    Charging a row to its start hour puts a 2 h window entirely in one bucket, which can
    push a single hour-of-day slot above the number of days in the range. Credited hours
    are spread over the row's clamped span pro rata.
    """
    mat = np.zeros((7, 24))
    if not len(df):
        return mat
    start = df.StartTime.values.astype("datetime64[s]").astype(np.int64)
    span = df.sec.values.astype(np.int64)
    val = df[col].values.astype(float)
    rate = np.divide(val, span, out=np.zeros(len(df)), where=span > 0)   # hours per second
    for t, remaining, r in zip(start, span, rate):
        if remaining <= 0 or r == 0:
            continue
        while remaining > 0:
            chunk = min(3600 - t % 3600, remaining)
            mat[(t // 86400 + 3) % 7, (t % 86400) // 3600] += chunk * r  # epoch day 0 = Thu
            t += chunk
            remaining -= chunk
    return mat

# ------------------------------------------------------------------ render
def txt(fig, x, y, s, size=12, color=INK, weight="normal", ha="left", va="baseline"):
    return fig.text(x, y, s, fontsize=size, color=color, fontweight=weight, ha=ha, va=va)

def footer(fig, x, y, s, size=11.0, maxw=.9163):
    """Footers carry the method notes, so they must not run off the canvas. Measure, then
    step the size down until the line fits the content width."""
    t = txt(fig, x, y, s, size, MUTED)
    fig.canvas.draw()
    while t.get_window_extent().width / fig.bbox.width > maxw and size > 10.0:
        size -= .25
        t.set_fontsize(size)
        fig.canvas.draw()
    if t.get_window_extent().width / fig.bbox.width > maxw:
        print("  ! footer still overflows at %.2fpt: %s..." % (size, s[:60]))
    return t

def section(fig, x, y, title, sub):
    txt(fig, x, y, title, 15.5, INK, "bold")
    txt(fig, x, y - .0195, sub, 11.5, MUTED)

def hbars(fig, rect, series, cmap, slots=10):
    ax = fig.add_axes(rect)
    slots = max(slots, len(series))
    y = np.arange(slots - len(series), slots)[::-1]
    vmax = series.max()
    ax.barh(y, series.values, .52, zorder=3,
            color=[cmap(.32 + .58 * (v / vmax)) for v in series])
    for yi, v in zip(y, series.values):
        ax.text(v + vmax * .02, yi, "{:,.0f} h".format(v), va="center", fontsize=11.5, color=INK2)
    ax.set_yticks(y)
    ax.set_yticklabels(series.index, fontsize=12, color=INK)
    ax.set_xlim(0, vmax * 1.30)
    ax.set_ylim(-.7, slots - .3)
    ax.set_xticks([])
    ax.set_facecolor(BG)
    for s in ax.spines.values():
        s.set_visible(False)
    ax.tick_params(length=0)

def dashboard(out, title, subtitle, kpis, months, stack, monthly_sub, cats, apps, heat,
              cmap, foot1, foot2, heat_sub, cat_sub, app_sub, annotate=True,
              xlabels=None, panel_title="Active hours per month", legend_cols=1,
              xlabel=None, legend_loc="upper left", head=1.0, note=None):
    fig = plt.figure(figsize=(19.2, 12.4), dpi=100)
    txt(fig, .0417, .955, title, 31, INK, "bold", va="center")
    txt(fig, .0417, .9275, subtitle, 13.5, INK2, va="center")

    cw, gap = .1760, .0125
    for i, (lab, val, unit) in enumerate(kpis):
        x = .0417 + i * (cw + gap)
        fig.patches.append(FancyBboxPatch((x, .7734), cw, .1129, transform=fig.transFigure,
            boxstyle="round,pad=0,rounding_size=0.004", linewidth=1,
            edgecolor=LINE, facecolor=CARD, zorder=0))
        txt(fig, x + .0198, .8605, lab, 12.5, INK2, va="center")
        t = txt(fig, x + .0198, .8125, val, 33, INK, "bold", va="center")
        fig.canvas.draw()
        w = t.get_window_extent().width / fig.bbox.width
        txt(fig, x + .0198 + w + .006, .8085, unit, 13, MUTED, va="center")

    ax = fig.add_axes([.048, .4435, .599, .250])
    xs = np.arange(len(months))
    bottom = np.zeros(len(months))
    for lab, vals, color in stack:
        ax.bar(xs, vals, .58, bottom=bottom, color=color, label=lab, zorder=3,
               linewidth=1.6 if bottom.any() else 0, edgecolor=CARD)
        bottom = bottom + np.asarray(vals, dtype=float)
    if annotate:
        k = int(np.argmax(bottom))
        ax.annotate("{:.0f} h".format(bottom[k]), (k, bottom[k]), textcoords="offset points",
                    xytext=(0, 9), ha="center", fontsize=12.5, color=INK, fontweight="bold")
    ax.set_xticks(xs)
    ax.set_xticklabels(xlabels if xlabels is not None else
                       ["%s\n%s" % (p.strftime("%b"), p.strftime("%y")) for p in months],
                       fontsize=11, color=INK2)
    ax.set_ylabel("hours", fontsize=10.5, color=MUTED)
    if xlabel:
        ax.set_xlabel(xlabel, fontsize=10.5, color=MUTED, labelpad=6)
    ax.grid(axis="y", color=LINE, linewidth=1, zorder=0)
    ax.set_axisbelow(True)
    for s in ("top", "right", "left"):
        ax.spines[s].set_visible(False)
    ax.spines["bottom"].set_color(LINE)
    ax.tick_params(length=0, colors=INK2, labelsize=11)
    if head > 1.0:
        ax.set_ylim(0, bottom.max() * head)          # headroom so the legend clears the bars
    if note:
        ax.annotate(note[0], xy=(note[1], bottom[note[1]]), xycoords="data",
                    xytext=(note[2], note[3]), textcoords="axes fraction",
                    fontsize=11.5, color=INK2, ha="left", va="center",
                    arrowprops=dict(arrowstyle="-", color=MUTED, linewidth=1,
                                    connectionstyle="angle3,angleA=0,angleB=70"))
    if len(stack) > 1:
        leg = ax.legend(frameon=False, fontsize=11.5, loc=legend_loc, handlelength=.9,
                        handleheight=.9, borderpad=0, labelspacing=.35,
                        ncol=legend_cols, columnspacing=1.6)
        for t in leg.get_texts():
            t.set_color(INK2)
    section(fig, .0417, .7305, panel_title, monthly_sub)

    hbars(fig, [.6958, .4470, .2050, .2465], cats, cmap)
    section(fig, .6958, .7305, "Where the time goes", cat_sub)

    ax = fig.add_axes([.0417, .0742, .606, .250])
    im = ax.imshow(heat, aspect="auto", cmap=cmap, vmin=0)
    ax.set_xticks(range(24))
    ax.set_xticklabels(["%02d" % h for h in range(24)], fontsize=10.5, color=INK2)
    ax.set_yticks(range(7))
    ax.set_yticklabels(["Mon", "Tue", "Wed", "Thu", "Fri", "Sat", "Sun"], fontsize=11.5, color=INK)
    ax.set_xlabel("hour of day", fontsize=10.5, color=MUTED, labelpad=6)
    ax.set_xticks(np.arange(-.5, 24, 1), minor=True)
    ax.set_yticks(np.arange(-.5, 7, 1), minor=True)
    ax.grid(which="minor", color=BG, linewidth=2)
    ax.tick_params(which="both", length=0)
    for s in ax.spines.values():
        s.set_visible(False)
    r, c = np.unravel_index(heat.argmax(), heat.shape)
    ax.text(c, r, "{:.0f} h".format(heat[r, c]), ha="center", va="center",
            fontsize=11.5, color="white", fontweight="bold")
    cax = fig.add_axes([.523, .3525, .124, .0075])
    cb = fig.colorbar(im, cax=cax, orientation="horizontal")
    cb.outline.set_visible(False)
    cax.tick_params(length=0, labelsize=10, colors=MUTED, pad=3)
    fig.text(.516, .3520, "hours", fontsize=10.5, color=MUTED, ha="right")
    section(fig, .0417, .3680, "Weekly rhythm", heat_sub)

    hbars(fig, [.6958, .0770, .2050, .2470], apps, cmap)
    section(fig, .6958, .3680, "Top applications", app_sub)

    footer(fig, .0417, .0300, foot1)
    footer(fig, .0417, .0097, foot2)
    fig.savefig(out, dpi=100, facecolor=BG)
    plt.close(fig)
    return out

def weekly_dashboard(out, a, mirrored_total, net_total):
    """Week-by-week slide: stacked category bars over a table of weekly totals."""
    w = a.copy()
    w["week"] = w.StartTime.dt.to_period("W-SUN")          # weeks run Mon-Sun
    weeks = pd.period_range(w.week.min(), w.week.max(), freq="W-SUN")
    w["cat"] = np.where(w.Category.isin(CAT_ORDER[:-1]), w.Category, "Other categories")
    stack = (w.pivot_table(index="week", columns="cat", values="net", aggfunc="sum")
              .reindex(index=weeks, columns=CAT_ORDER).fillna(0))
    totals = stack.sum(axis=1)
    peak = totals.idxmax()

    fig = plt.figure(figsize=(19.2, 12.4), dpi=100)
    txt(fig, .0417, .955, "Computer Activity — week by week, all PCs", 31, INK, "bold", va="center")
    txt(fig, .0417, .9275, "all 4 PCs summed, with remote views and simultaneous use counted once  ·  "
        "%s – %s  ·  %d weeks" % (fmt_date(a.StartTime.min()),
                                            fmt_date(a.StartTime.max()), len(weeks)),
        13.5, INK2, va="center")

    kpis = [("Active hours logged", "{:,.0f}".format(net_total), "h"),
            ("Weeks in the log", str(len(weeks)), "weeks"),
            ("Avg per week", "{:.1f}".format(totals.mean()), "h"),
            ("Busiest week", "{:.0f}".format(totals.max()),
             "h  w/c %d %s" % (peak.start_time.day, peak.start_time.strftime("%b %y"))),
            ("Double counting removed", "-{:,.0f}".format(mirrored_total), "h")]
    cw, gap = .1760, .0125
    for i, (lab, val, unit) in enumerate(kpis):
        x = .0417 + i * (cw + gap)
        fig.patches.append(FancyBboxPatch((x, .7734), cw, .1129, transform=fig.transFigure,
            boxstyle="round,pad=0,rounding_size=0.004", linewidth=1,
            edgecolor=LINE, facecolor=CARD, zorder=0))
        txt(fig, x + .0198, .8605, lab, 12.5, INK2, va="center")
        t = txt(fig, x + .0198, .8125, val, 33, INK, "bold", va="center")
        fig.canvas.draw()
        wd = t.get_window_extent().width / fig.bbox.width
        txt(fig, x + .0198 + wd + .006, .8085, unit, 13, MUTED, va="center")

    # ---- stacked bars, one per week
    ax = fig.add_axes([.048, .4560, .910, .2390])
    xs = np.arange(len(weeks))
    bottom = np.zeros(len(weeks))
    for cat, color in CAT_COLORS:
        v = stack[cat].values
        ax.bar(xs, v, .70, bottom=bottom, color=color, zorder=3, linewidth=.9, edgecolor=CARD,
               label="%s  (%s h)" % (cat, format(int(round(v.sum())), ",")))
        bottom = bottom + v
    k = int(np.argmax(bottom))
    ax.annotate("{:.0f} h".format(bottom[k]), (k, bottom[k]), textcoords="offset points",
                xytext=(0, 8), ha="center", fontsize=12.5, color=INK, fontweight="bold")
    firsts = [i for i, p in enumerate(weeks)
              if p.start_time.day <= 7 and p.start_time.month % 2 == 1]
    ax.set_xticks(firsts)
    ax.set_xticklabels(["%s\n%s" % (weeks[i].start_time.strftime("%b"),
                                    weeks[i].start_time.strftime("%y")) for i in firsts],
                       fontsize=11, color=INK2)
    ax.set_xlim(-.8, len(weeks) - .2)
    ax.set_ylabel("hours", fontsize=10.5, color=MUTED)
    ax.grid(axis="y", color=LINE, linewidth=1, zorder=0)
    ax.set_axisbelow(True)
    for s in ("top", "right", "left"):
        ax.spines[s].set_visible(False)
    ax.spines["bottom"].set_color(LINE)
    ax.tick_params(length=0, colors=INK2, labelsize=11)
    leg = ax.legend(frameon=False, fontsize=11.5, loc="upper left", handlelength=.9,
                    handleheight=.9, borderpad=0, labelspacing=.3, ncol=3, columnspacing=1.6)
    for t in leg.get_texts():
        t.set_color(INK2)
    section(fig, .0417, .7305, "Where the time goes, week by week",
            "active hours per week, stacked by category — all PCs, no second counted twice")

    # ---- table of weekly totals
    section(fig, .0417, .3760, "Total active hours per week",
            "every week in the log; the tint behind each number is scaled to the busiest week")
    nblocks = 4
    per = int(np.ceil(len(weeks) / nblocks))
    bw, bgap = .2130, .0230
    top_y, rh = .3240, .0176
    vmax = totals.max()
    for b in range(nblocks):
        x0 = .0417 + b * (bw + bgap)
        txt(fig, x0, top_y + .009, "WEEK OF", 10, MUTED)
        txt(fig, x0 + bw, top_y + .009, "HOURS", 10, MUTED, ha="right")
        fig.patches.append(plt.Rectangle((x0, top_y + .0035), bw, .0011,
                                         transform=fig.transFigure, color=LINE, zorder=1))
        for r in range(per):
            i = b * per + r
            if i >= len(weeks):
                break
            y = top_y - (r + 1) * rh
            v = totals.iloc[i]
            fig.patches.append(plt.Rectangle((x0, y + .0016), bw * (v / vmax), rh - .004,
                                             transform=fig.transFigure, zorder=1,
                                             color=mix("#2A78D6", "white", .88)))
            d = weeks[i].start_time
            bold = "bold" if i == int(np.argmax(totals.values)) else "normal"
            txt(fig, x0 + .004, y + .0055, "%d %s %s" % (d.day, d.strftime("%b"), d.strftime("%y")),
                11.5, INK if bold == "bold" else INK2, bold)
            txt(fig, x0 + bw - .004, y + .0055, "{:,.1f}".format(v), 11.5, INK, bold, ha="right")

    footer(fig, .0417, .0300,
        "Source: ActivityLogger-data/data/activity_slim.csv.gz (3 work PCs) + the 2 local CSVs for "
        "DESKTOP-GHT99IR. Weeks run Monday to Sunday; the first and last weeks are partial.")
    footer(fig, .0417, .0097,
        "Excludes idle, the lock screen and windows held over 2 h by a scheduled task. RealVNC time is dropped "
        "— always a view of one of these PCs — and simultaneous use counted once: {:,.0f} h in all. Categories "
        "outside the top five are “Other”.".format(mirrored_total))

    fig.savefig(out, dpi=100, facecolor=BG)
    plt.close(fig)
    return out, totals, peak

# ------------------------------------------------------------------ shape the data
def panels(a, hours_col="hours"):
    cats = a.groupby("Category")[hours_col].sum().sort_values(ascending=False)
    small = cats[cats < max(5, cats.sum() * .004)]
    if len(small) > 1:
        cats = pd.concat([cats[~cats.index.isin(small.index)],
                          pd.Series({"Other categories": small.sum()})])
    cats = cats.head(10)
    apps = (a.groupby(a.ProcessName.str.replace(r"\.exe$", "", regex=True, case=False))
             [hours_col].sum().sort_values(ascending=False).head(10))
    return cats, apps, hour_split(a, hours_col)

def fmt_date(ts):
    return "%d %s %d" % (ts.day, ts.strftime("%b"), ts.year)

def month_index(a):
    return pd.period_range(a.StartTime.min().to_period("M"), a.StartTime.max().to_period("M"),
                           freq="M")

def by_month(a, idx, col="hours", mask=None):
    g = (a[mask] if mask is not None else a)
    return g.groupby(g.StartTime.dt.to_period("M"))[col].sum().reindex(idx, fill_value=0).values

# weeks that get their own hour-by-hour slide
WEEKS = [pd.Timestamp("2026-03-16"), pd.Timestamp("2026-05-18")]

def ordinal(n):
    return "%d%s" % (n, "th" if 10 <= n % 100 <= 20 else
                     {1: "st", 2: "nd", 3: "rd"}.get(n % 10, "th"))

def week_slide(a, wtotals, w0):
    """One week, hour by hour, stacked by PC."""
    w1 = w0 + pd.Timedelta(days=7)
    rank = int((wtotals > wtotals.loc[w0.to_period("W-SUN")]).sum()) + 1
    wk = a[(a.StartTime >= w0) & (a.StartTime < w1)]
    wk_h = wk.net.sum()
    by_hour = {pc: hour_split(wk[wk.ComputerName == pc]).sum(axis=0) for pc, _, _ in PCS}
    hstack = [("%s  (%s h)" % (pc, format(round(by_hour[pc].sum(), 1), ",")),
               by_hour[pc], COLOR[pc]) for pc, _, _ in PCS if by_hour[pc].sum() > .05]
    tot_hour = sum(v for _, v, _ in hstack)
    cats, apps, heat = panels(wk, "net")
    peak_hr = int(np.argmax(tot_hour))
    night = tot_hour[:4].sum()
    span = "%s – %s" % (fmt_date(w0), fmt_date(w1 - pd.Timedelta(days=1)))
    print("week of %s: %.1f h, rank %d of %d, peak %.1f h at %02d:00, %.1f h in 00:00-04:00"
          % (w0.date(), wk_h, rank, len(wtotals), tot_hour[peak_hr], peak_hr, night))
    png = dashboard(
        os.path.join(REPO, "ActivityDashboard-Week-%s.png" % w0.strftime("%Y-%m-%d")),
        "Computer Activity — week of %d %s" % (w0.day, w0.strftime("%B %Y")),
        "%s  ·  %s active hours across %d PCs, counted once  ·  the %s busiest week of %d"
        % (span, format(round(wk_h, 1), ","), len(hstack), ordinal(rank), len(wtotals)),
        [("Active hours", "{:.0f}".format(wk_h), "h"),
         ("Days with activity", str(wk.loc[wk.net > 0, "StartTime"].dt.date.nunique()), "days"),
         ("Avg per day", "{:.1f}".format(wk_h / 7), "h"),
         ("Busiest hour", "{:.1f}".format(tot_hour[peak_hr]), "h  at %02d:00" % peak_hr),
         ("vs. weekly average", "{:+.0f}".format(100 * wk_h / wtotals.mean() - 100), "%")],
        np.arange(24), hstack,
        "every row split across the clock hours it actually covers, then stacked by PC",
        cats, apps, heat, BLUE_RAMP,
        "Source: the same rows as the week-by-week slide, filtered to %s. RealVNC time is dropped and "
        "simultaneous use counted once, so this sums to the %s h that week shows."
        % (span, format(round(wk_h, 1), ",")),
        "Windows held over 2 h are excluded, so the scheduled task that woke HC-D2BYW34 around 01:38 does not "
        "read as a night shift: {:.1f} h remains between 00:00 and 04:00. Excludes idle and the lock screen."
        .format(night),
        "the same hours, split by day of that week", "active hours by category, that week",
        "active hours by process (.exe omitted), that week",
        xlabels=["%02d" % h for h in range(24)], panel_title="Active hours by hour of day",
        annotate=False, legend_cols=2, xlabel="hour of day", legend_loc="upper right", head=1.28)
    return (png, "Week of %d %s - %s active hours, hour by hour and stacked by PC"
            % (w0.day, w0.strftime("%B %Y"), format(round(wk_h, 1), ",")))

def main():
    df, n_local_raw, n_local, n_work = load()
    a, idle_h, n_failed, failed_h = active(df)
    print("dropped %d idle-detection failures worth %.0f h" % (n_failed, failed_h))
    a["net"] = credited_seconds(a) / 3600
    a["mirror_h"] = a.hours - a.net
    mirrored_total = a.mirror_h.sum()
    vnc_removed = a.loc[a.viewer, "mirror_h"].sum()
    print("PC active hours (gross):", a.groupby("ComputerName").hours.sum().round(1).to_dict())
    print("credited:", a.groupby("ComputerName").net.sum().round(1).to_dict())
    print("removed as simultaneous:", a.groupby("ComputerName").mirror_h.sum().round(1).to_dict())
    print("gross %.1f  net %.1f  removed %.1f (of which VNC %.1f)"
          % (a.hours.sum(), a.net.sum(), mirrored_total, vnc_removed))

    idx = month_index(a)
    days = a.StartTime.dt.date.nunique()
    slides = []

    # ---------------- overview
    cats, apps, heat = panels(a, "net")
    stack = []
    for pc, nice, hue in PCS:
        g = a[a.ComputerName == pc]
        stack.append(("%s \u2014 %s  (%s h)" % (pc, nice, format(int(round(g.net.sum())), ",")),
                      by_month(g, idx, "net"), hue))
    net_total = a.net.sum()
    slides.append((dashboard(
        os.path.join(REPO, "ActivityDashboard-AllPCs.png"),
        "Computer Activity \u2014 all PCs",
        "ActivityLogger foreground-window logs from 4 PCs  \u00b7  %s \u2013 %s  \u00b7  "
        "%d days with recorded activity" % (fmt_date(a.StartTime.min()),
                                            fmt_date(a.StartTime.max()), days),
        [("Active hours logged", "{:,.0f}".format(net_total), "h"),
         ("Days with activity", str(days), "days"),
         ("Avg per active day", "{:.1f}".format(net_total / days), "h"),
         ("PCs logged", "4", "machines"),
         ("Double counting removed", "-{:,.0f}".format(mirrored_total), "h")],
        idx, stack,
        "each PC's own log, minus every remote view of another PC and any second already counted elsewhere",
        cats, apps, heat, BLUE_RAMP,
        "Source: ActivityLogger-data/data/activity_slim.csv.gz (3 work PCs, 305,961 deduped rows) "
        "+ %s local CSVs for DESKTOP-GHT99IR (%s raw \u2192 %s deduped). PCs grouped by the ComputerName column."
        % (2, format(n_local_raw, ","), format(n_local, ",")),
        "Excludes idle, the lock screen and {:,.0f} windows held over 2 h where idle detection failed ({:,.0f} h). All "
        "{:,.0f} h of RealVNC is dropped and {:,.0f} h of two-PC overlap counted once. Jun 25 and Aug 26 are partial."
        .format(n_failed, failed_h, vnc_removed, mirrored_total - vnc_removed),
        "total active hours by weekday and hour of day, all PCs combined",
        "active hours by category, all PCs combined",
        "active hours by process (.exe omitted), all PCs combined"),
        "All PCs - %s to %s - %s active hours across 4 machines, double counting removed (-%s h)"
        % (fmt_date(a.StartTime.min()), fmt_date(a.StartTime.max()),
           format(int(round(net_total)), ","), format(int(round(mirrored_total)), ","))))

    # ---------------- week by week, all PCs
    wpng, wtotals, wpeak = weekly_dashboard(
        os.path.join(REPO, "ActivityDashboard-Weekly.png"), a, mirrored_total, net_total)
    slides.append((wpng,
        "Week by week, all PCs - %d weeks, %s active hours, busiest week %s h w/c %d %s"
        % (len(wtotals), format(int(round(net_total)), ","), format(round(wtotals.max(), 1), ","),
           wpeak.start_time.day, wpeak.start_time.strftime("%b %Y"))))

    # ---------------- one hour-by-hour slide per week of interest, oldest first
    for w0 in sorted(WEEKS):
        slides.append(week_slide(a, wtotals, w0))

    # ---------------- one slide per PC, largest first
    order = a.groupby("ComputerName").hours.sum().sort_values(ascending=False).index
    for pc in order:
        g = a[a.ComputerName == pc].copy()
        hue = COLOR[pc]
        cmap = ramp_for(hue)
        cats, apps, heat = panels(g)
        gidx = month_index(g)
        gh, gdays = g.hours.sum(), g.StartTime.dt.date.nunique()
        vh, mh = g[g.viewer].hours.sum(), g.mirror_h.sum()
        vmh = g.loc[g.viewer, "mirror_h"].sum()        # the VNC-only slice of that overlap
        gidle = df[(df.ComputerName == pc) &
                   ((df.Category == "Inactive") | (df.ProcessName.str.lower() == "lockapp.exe"))
                   ].DurationSeconds.clip(upper=7200).sum() / 3600
        if vh > gh * .02:
            stack = [("Local apps  ({:,.0f} h)".format(gh - vh),
                      by_month(g, gidx, "hours", ~g.viewer), hue),
                     ("Remote session (VNC)  ({:,.0f} h — a view of another PC, counted there instead)"
                      .format(vh), by_month(g, gidx, "hours", g.viewer), mix(hue, "white", .5))]
            msub = "split by whether the foreground window was a local app or a RealVNC view of another PC"
        else:
            stack = [("", by_month(g, gidx, "hours"), hue)]
            msub = "hours this PC logged itself"
        kpis = [("Active hours logged", "{:,.0f}".format(gh), "h"),
                ("Days with activity", str(gdays), "days"),
                ("Avg per active day", "{:.1f}".format(gh / gdays), "h"),
                ("Window switches", "{:,.1f}K".format(len(g) / 1000), "rows"),
                ("Excluded from the\nall-PC totals", "{:,.0f}".format(mh), "h")]
        local = pc == "DESKTOP-GHT99IR"
        slides.append((dashboard(
            os.path.join(REPO, "ActivityDashboard-%s.png" % pc),
            "Computer Activity \u2014 %s" % pc,
            "%s  \u00b7  %s \u2013 %s  \u00b7  %d days with recorded activity"
            % ("This machine" if local else NAME[pc],
               fmt_date(g.StartTime.min()), fmt_date(g.StartTime.max()), gdays),
            kpis, gidx, stack, msub, cats, apps, heat, cmap,
            "Source: %s" % ("the 2 local CSVs in .../OneDrive/Documents/ActivityLogger "
                            "(%s raw rows \u2192 %s after removing the overlapping snapshot copy)."
                            % (format(n_local_raw, ","), format(n_local, ","))
                            if local else
                            "ActivityLogger-data/data/activity_slim.csv.gz, rows where "
                            "ComputerName = %s. Window titles were stripped from that export." % pc),
            "Excludes idle, the lock screen (%s h here) and windows held over 2 h by a scheduled task. This is the PC's "
            "own log in full; the all-PC slides leave out %s h of it%s."
            % (format(int(round(gidle)), ","), format(int(round(mh)), ","),
               " — the %s h of RealVNC views of the other machines, plus overlap"
               % format(int(round(vh)), ",") if vh > gh * .02
               else ", already counted on another PC"),
            "total active hours by weekday and hour of day",
            "active hours by category", "active hours by process (.exe omitted)"),
            "%s - %s - %s to %s - %s active hours%s"
            % (pc, NAME[pc], fmt_date(g.StartTime.min()), fmt_date(g.StartTime.max()),
               format(int(round(gh)), ","),
               ", plus %s h of RealVNC views counted on the PC being driven" % format(int(round(vh)), ",")
               if vh > gh * .02 else "")))
    return slides, a, mirrored_total, net_total, days

# ------------------------------------------------------------------ deck
def build_deck(slides):
    """Rebuild the pptx: 13.33x7.5 blank slides, picture box and caption as before."""
    prs = Presentation(DECK) if os.path.exists(DECK) else Presentation()
    ref_bodypr = None
    if len(prs.slides) and len(prs.slides[0].shapes) >= 2:
        ref_pic, ref_cap = prs.slides[0].shapes[0], prs.slides[0].shapes[1]
        ref_bodypr = copy.deepcopy(ref_cap.text_frame._txBody.find(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr"))
        geom = (ref_pic.left, ref_pic.top, ref_pic.width, ref_pic.height,
                ref_cap.left, ref_cap.top, ref_cap.width, ref_cap.height)
    else:                                                 # first run: 16:9, same margins
        prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
        geom = tuple(Inches(v) for v in (1.433, .22, 10.467, 6.76, 1.433, 7.04, 10.467, .3))

    xml_slides = prs.slides._sldIdLst                     # drop the old slides
    for sld in list(xml_slides):
        rId = sld.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
        prs.part.drop_rel(rId)
        xml_slides.remove(sld)

    for png, caption in slides:
        s = prs.slides.add_slide(prs.slide_layouts[6])
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = RGBColor.from_string("F9F9F7")
        s.shapes.add_picture(png, geom[0], geom[1], geom[2], geom[3])
        box = s.shapes.add_textbox(geom[4], geom[5], geom[6], geom[7])
        if ref_bodypr is not None:                        # keep the original wrap/autofit
            tx = box.text_frame._txBody
            tx.remove(tx.find("{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr"))
            tx.insert(0, copy.deepcopy(ref_bodypr))
        p = box.text_frame.paragraphs[0]
        p.font.size = Pt(11)
        p.font.name = "Segoe UI"
        p.font.color.rgb = RGBColor.from_string("898781")
        p.add_run().text = caption
    prs.save(DECK)

if __name__ == "__main__":
    slides, a, mirrored, net_total, days = main()
    for p, cap in slides:
        print("wrote", os.path.basename(p))
    build_deck(slides)
    chk = Presentation(DECK)
    print("\ndeck: %d slides" % len(chk.slides))
    for i, s in enumerate(chk.slides):
        print(" ", i + 1, [sh.text_frame.text for sh in s.shapes if sh.has_text_frame][0])
